from pptx import Presentation
from pathlib import Path
from my_argparse import get_args
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR, MSO_AUTO_SIZE
from pptx.dml.color import RGBColor
import os

TEMPLATE_FILENAME = "Template.pptx"
COMMENT_NOTATION = "#"
BLANK_SLIDE_NOTATION = "!"
FONT = "ARIAL"
TITLE_FONT_SIZE = 40
LYRICS_FONT_SIZE = 27

def get_template_blank_layout():
    template_ppt = Presentation(input_file_path(TEMPLATE_FILENAME))
    blank_layout = template_ppt.slide_layouts[0]
    return blank_layout


def data_folder_path():
    return Path(__file__).parent.parent.joinpath("data")


def input_folder_path():
    return f"{data_folder_path()}/input/"


def input_file_path(input_filename):
    return f"{data_folder_path()}/input/{input_filename}"


def output_file_path():
    return f"{data_folder_path()}/output"


def read_file_into_blocks(input_file, make_uppercase):
    file_contents = ""
    with open(input_file, "r") as file:
        for line in file:
            if not line.startswith(COMMENT_NOTATION):
                file_contents += line
    if make_uppercase:
        file_contents = file_contents.upper()
    return file_contents.split("\n\n")


def create_title_slide(ppt, blank_layout, title, authors):
    title_slide = ppt.slides.add_slide(blank_layout)
    create_textbox(ppt, title_slide, f"{title}\n{authors}")


def create_blank_slide(ppt, blank_layout):
    ppt.slides.add_slide(blank_layout)


def create_lyric_slide(ppt, blank_layout, lyric_block):
    lyric_slide = ppt.slides.add_slide(blank_layout)
    create_textbox(ppt, lyric_slide, lyric_block)


def create_textbox(ppt, slide, contents):
    txBox = slide.shapes.add_textbox(
        0, 0, ppt.slide_width, ppt.slide_height)
    
    tf = txBox.text_frame
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE

    p = tf.paragraphs[0]
    p.text = contents
    p.alignment = PP_ALIGN.CENTER

    font = p.font
    font.name = FONT
    font.size = Pt(LYRICS_FONT_SIZE)
    font.color.rgb = RGBColor(255, 255, 255)
    font.bold = True


def create_blank_layout(ppt):
    blank_layout = ppt.slide_layouts[6]
    fill = blank_layout.background.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(0, 0, 0)


def create_ppt(title, authors, lyric_blocks):
    ppt = Presentation()
    blank_layout = create_blank_layout(ppt)
    
    blank_layout = ppt.slide_layouts[6]
    fill = blank_layout.background.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(0, 0, 0)

    create_title_slide(ppt, blank_layout, title, authors)
    for lyric_block in lyric_blocks:
        if lyric_block.startswith(BLANK_SLIDE_NOTATION):
            create_blank_slide(ppt, blank_layout)
        else:
            create_lyric_slide(ppt, blank_layout, lyric_block)
    return ppt


def get_title_and_authors(blocks):
    return blocks[0], blocks[1]


def get_lyric_blocks(blocks):
    return blocks[2:]


def process_lyrics_text_file(filename, make_uppercase):
    print(f"Reading file {filename}...")
    blocks = read_file_into_blocks(filename, make_uppercase)
    title, authors = get_title_and_authors(blocks)
    
    print(f"Making slides for {title} by {authors}...")
    lyric_blocks = get_lyric_blocks(blocks)
    ppt = create_ppt(title, authors, lyric_blocks)
    
    print(f"Saving file {title}.pptx to {output_file_path()}...")
    ppt.save(f"{output_file_path()}/{title}.pptx")


def main():
    args = get_args()
    make_uppercase = args.uppercase
    for file in os.listdir(input_folder_path()):
        filename = input_folder_path() + os.fsdecode(file)
        if filename.endswith("txt"):
            process_lyrics_text_file(filename, make_uppercase)


if __name__ == "__main__":
    main()
